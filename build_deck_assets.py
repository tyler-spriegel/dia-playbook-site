#!/usr/bin/env python3
"""Extract table shapes from the DIA playbook pptx into deck-tables.js."""
from __future__ import annotations

import html
import json
from pathlib import Path

from pptx import Presentation

PPTX = Path(__file__).resolve().parent.parent / "[DRAFT] DIA Delivery Playbook.pptx"
OUT = Path(__file__).resolve().parent / "deck-tables.js"

# PowerPoint for the web (SharePoint). wdSlideIndex is 0-based (slide 1 → 0).
PLAYBOOK_DECK_ONLINE_URL = (
    "https://pwc.sharepoint.com/:p:/r/teams/US-INT-ADVS-Enterprise-Application-Integration/"
    "Shared%20Documents/10.%20Large%20Scale%20Solution%20delivery/2.%20Delivery/"
    "Delivery%20Playbook/%5BDRAFT%5D%20DIA%20Delivery%20Playbook.pptx"
    "?d=w5cbc36da055d4fc9af36433a8a8d0b38&csf=1&web=1&e=dDRHTh"
)


def deck_slide_url(slide_1based: int) -> str:
    return f"{PLAYBOOK_DECK_ONLINE_URL}&wdSlideIndex={slide_1based - 1}"


def deck_slide_href(slide_1based: int) -> str:
    return html.escape(deck_slide_url(slide_1based), quote=True)


def slide_anchor(slide_1based: int, label: str) -> str:
    h = deck_slide_href(slide_1based)
    return (
        f'<a class="deck-slide-link" href="{h}" target="_blank" rel="noopener noreferrer">'
        f"{html.escape(label)}</a>"
    )


def slide_index_links(lo: int, hi: int) -> str:
    return " · ".join(slide_anchor(i, str(i)) for i in range(lo, hi + 1))

# section_id -> list of { key, slide, title }
SECTION_TABLES: dict[str, list[dict]] = {
    "intro": [
        {"key": "toc", "slide": 7, "title": "Contents (Slide 7)"},
    ],
    "operating-model": [
        {"key": "raci_program", "slide": 15, "title": "Sample RACI — Program workstreams (Slide 15)"},
        {"key": "pm_roles", "slide": 11, "title": "Roles — Project Management (Slide 11)"},
    ],
    "mobilization": [],
    "agile": [
        {"key": "sprint_calendar", "slide": 34, "title": "Program sprint cadence — calendar (Slide 34)"},
    ],
    "testing": [
        {"key": "uat_roles", "slide": 49, "title": "UAT expectations by role (Slide 49)"},
        {"key": "uat_criteria", "slide": 50, "title": "UAT entry & exit criteria (Slide 50)"},
    ],
    "release": [],
    "business-readiness": [
        {"key": "training_methods", "slide": 74, "title": "Training development methods (Slide 74)"},
    ],
    "golive": [
        {"key": "gonogo_outcomes", "slide": 85, "title": "Go / No-Go meeting outcomes (Slide 85)"},
    ],
    "hypercare": [
        {"key": "defect_sla", "slide": 89, "title": "Hypercare defect management (Slide 89, table 1)"},
    ],
    "tooling": [
        {"key": "tool_roles", "slide": 91, "title": "Tooling by role (Slide 91)"},
    ],
    "appendix": [
        {"key": "testing_cycles_a", "slide": 97, "title": "Testing cycles — definitions (Slide 97)"},
        {"key": "raci_workstream", "slide": 103, "title": "Sample RACI — Workstream view (Slide 103)"},
    ],
}


def extract_tables(prs: Presentation, slide_idx: int) -> list[list[list[str]]]:
    slide = prs.slides[slide_idx - 1]
    out: list[list[list[str]]] = []
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        rows = []
        for r in range(len(tbl.rows)):
            row = []
            for c in range(len(tbl.columns)):
                cell = tbl.cell(r, c)
                text = "\n".join(
                    p.text for p in cell.text_frame.paragraphs if p.text
                )
                text = (
                    text.replace("\r", " ")
                    .replace("\u200b", "")
                    .replace("\xa0", " ")
                    .strip()
                )
                row.append(text)
            rows.append(row)
        out.append(rows)
    return out


def raci_rowspan_html(rows: list[list[str]]) -> str:
    """Program RACI: col0 workstream (merged), col1 activity, cols 2–6 R/A/C/I letters."""
    if not rows:
        return ""
    header = rows[0]
    body = rows[1:]
    parts = ['<thead><tr>']
    for h in header:
        parts.append(f"<th>{html.escape(h)}</th>")
    parts.append("</tr></thead><tbody>")
    i = 0
    while i < len(body):
        row = list(body[i])
        while len(row) < 7:
            row.append("")
        ws, activity = row[0], row[1]
        letters = row[2:7]
        if ws:
            span = 1
            j = i + 1
            while j < len(body):
                rj = list(body[j])
                while len(rj) < 7:
                    rj.append("")
                if rj[0]:
                    break
                span += 1
                j += 1
            parts.append("<tr>")
            if span > 1:
                parts.append(
                    f'<th scope="row" rowspan="{span}" class="raci-ws">{html.escape(ws)}</th>'
                )
            else:
                parts.append(
                    f'<th scope="row" class="raci-ws">{html.escape(ws)}</th>'
                )
            parts.append(
                f'<td class="raci-activity">{html.escape(activity).replace(chr(10), "<br/>")}</td>'
            )
            for c in letters:
                parts.append(raci_cell(c))
            parts.append("</tr>")
            for k in range(i + 1, j):
                r2 = list(body[k])
                while len(r2) < 7:
                    r2.append("")
                act2 = r2[1]
                let2 = r2[2:7]
                parts.append("<tr>")
                parts.append(
                    f'<td class="raci-activity">{html.escape(act2).replace(chr(10), "<br/>")}</td>'
                )
                for c in let2:
                    parts.append(raci_cell(c))
                parts.append("</tr>")
            i = j
        else:
            i += 1
    parts.append("</tbody>")
    return "".join(parts)


def raci_cell(text: str) -> str:
    t = html.escape(text)
    cls = "raci-letter"
    if text.strip() in ("A", "R", "C", "I"):
        cls += f" raci-{text.strip().lower()}"
    elif "A/R" in text or text == "A/R":
        cls += " raci-ar"
    return f'<td class="{cls}">{t.replace(chr(10), "<br/>")}</td>'


def simple_table_html(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    parts = ["<thead><tr>"]
    for h in rows[0]:
        parts.append(f"<th>{html.escape(h)}</th>")
    parts.append("</tr></thead><tbody>")
    for row in rows[1:]:
        parts.append("<tr>")
        for cell in row:
            parts.append(
                f'<td>{html.escape(cell).replace(chr(10), "<br/>")}</td>'
            )
        parts.append("</tr>")
    parts.append("</tbody>")
    return "".join(parts)


def table_to_html(rows: list[list[str]], mode: str) -> str:
    if mode == "raci_program" and rows and len(rows[0]) >= 5:
        return raci_rowspan_html(rows)
    return simple_table_html(rows)


def wrap_figure(
    inner_table: str, slide: int, title: str, note: str | None = None
) -> str:
    open_link = slide_anchor(slide, "Open this slide in deck")
    note = note or (
        "Table text extracted from the PowerPoint file. "
        "Use the link above for the full slide layout, graphics, and any content outside this table."
    )
    return (
        f'<figure class="slide-capture" data-slide="{slide}">'
        f'<figcaption class="slide-capture__title">{html.escape(title)} · {open_link}</figcaption>'
        f'<div class="slide-capture__frame">'
        f'<div class="slide-capture__chrome" aria-hidden="true"></div>'
        f'<div class="slide-capture__body"><div class="slide-capture__scroll">'
        f'<table class="slide-deck-table data-table">{inner_table}</table>'
        f"</div></div></div>"
        f'<p class="slide-capture__note">{html.escape(note)}</p>'
        f"</figure>"
    )


def main() -> None:
    prs = Presentation(str(PPTX))
    nslides = len(prs.slides)
    by_section: dict[str, list[str]] = {k: [] for k in SECTION_TABLES}

    for section, items in SECTION_TABLES.items():
        for item in items:
            slide_n = item["slide"]
            if slide_n < 1 or slide_n > nslides:
                continue
            tables = extract_tables(prs, slide_n)
            if not tables:
                continue
            tidx = 0
            rows = tables[tidx]
            if item["key"] == "raci_workstream" and len(rows) > 2:
                rows = [rows[0]] + [r for r in rows[1:] if any(c.strip() for c in r)]
                rows[0] = [
                    "ID",
                    "Core Program Governance Activities",
                    "Executive Sponsor",
                    "Program Management",
                    "Architecture",
                    "Delivery",
                    "Release Readiness",
                    "DIA Workstream",
                ]
            mode = item["key"]
            inner = table_to_html(rows, mode)
            if not inner:
                continue
            fig = wrap_figure(inner, slide_n, item["title"])
            by_section[section].append(fig)

    # JSON-escape each section's concatenated HTML
    payload = {}
    hints = [
        (
            "mobilization",
            "Mobilization & readiness in the SharePoint deck",
            "These slides are mostly text and diagrams (no table shapes were extracted). Jump to the matching slide numbers in the online file:",
            [(17, 24), (108, 110)],
        ),
        (
            "release",
            "Release management in the SharePoint deck",
            "Diagrams and narrative slides—open by number in the online file:",
            [(54, 61)],
        ),
        (
            "reporting",
            "Reporting & metrics in the SharePoint deck",
            "Section header slides in the draft—open by number in the online file:",
            [(75, 80)],
        ),
    ]
    for section, cap, intro, ranges in hints:
        blocks = []
        for lo, hi in ranges:
            label = f"Slides {lo}–{hi}" if lo != hi else f"Slide {lo}"
            blocks.append(
                f"<p><strong>{html.escape(label)}:</strong> {slide_index_links(lo, hi)}</p>"
            )
        ph = (
            f'<figure class="slide-capture slide-capture--deck-ref">'
            f'<figcaption class="slide-capture__title">{html.escape(cap)}</figcaption>'
            f'<div class="slide-capture__frame slide-capture__frame--deck-ref">'
            f'<p class="slide-capture__deck-intro">{html.escape(intro)}</p>'
            f'{"".join(blocks)}'
            f"</div></figure>"
        )
        by_section[section] = by_section.get(section, []) + [ph]

    for section, fragments in by_section.items():
        if not fragments:
            payload[section] = ""
        else:
            payload[section] = "\n".join(fragments)

    js = (
        "/* Auto-generated by build_deck_assets.py — do not edit by hand */\n"
        "window.PLAYBOOK_DECK_URL = "
        + json.dumps(PLAYBOOK_DECK_ONLINE_URL)
        + ";\n"
        "window.DECK_TABLES_BY_SECTION = "
        + json.dumps(payload, ensure_ascii=False)
        + ";\n"
    )
    OUT.write_text(js, encoding="utf-8")
    print(f"Wrote {OUT} ({len(js)} bytes)")


if __name__ == "__main__":
    main()
